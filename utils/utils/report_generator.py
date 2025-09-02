from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import pandas as pd
import plotly.graph_objects as go
import plotly.io as pio
import io
import base64
from typing import Dict, List
import tempfile
import os

class PowerPointReportGenerator:
    """Générateur de rapports PowerPoint automatisés"""
    
    def __init__(self):
        self.colors = {
            'primary': RGBColor(52, 152, 219),      # Bleu
            'secondary': RGBColor(46, 204, 113),     # Vert
            'accent': RGBColor(231, 76, 60),         # Rouge
            'warning': RGBColor(243, 156, 18),       # Orange
            'dark': RGBColor(52, 73, 94),            # Gris foncé
            'light': RGBColor(236, 240, 241)         # Gris clair
        }
        
    def create_presentation(self, parsed_data: Dict, figures: Dict) -> Presentation:
        """Crée une présentation PowerPoint complète"""
        prs = Presentation()
        
        # Slide 1: Page de titre
        self._create_title_slide(prs, parsed_data)
        
        # Slide 2: Résumé exécutif
        self._create_executive_summary_slide(prs, parsed_data)
        
        # Slide 3: Analyse mensuelle
        self._create_monthly_analysis_slide(prs, parsed_data, figures)
        
        # Slide 4: Performance des agents
        self._create_agents_performance_slide(prs, parsed_data, figures)
        
        # Slide 5: Tendances et KPIs
        self._create_kpi_trends_slide(prs, parsed_data, figures)
        
        # Slide 6: Analyse de la résolution
        self._create_resolution_analysis_slide(prs, parsed_data)
        
        # Slide 7: Recommandations
        self._create_recommendations_slide(prs, parsed_data)
        
        return prs
    
    def _create_title_slide(self, prs: Presentation, data: Dict):
        """Crée la slide de titre"""
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = "Rapport de Performance Téléphonie"
        title.text_frame.paragraphs[0].font.color.rgb = self.colors['primary']
        
        # Calculer la période couverte
        monthly_data = data.get('monthly_data', pd.DataFrame())
        if not monthly_data.empty and 'mois' in monthly_data:
            first_month = monthly_data['mois'].iloc[0]
            last_month = monthly_data['mois'].iloc[-1]
            period = f"Période: {first_month} - {last_month} 2025"
        else:
            period = "Période: 2025"
        
        subtitle.text = f"{period}\nAnalyse Automatisée"
        
        # Ajouter la date de génération
        date_shape = slide.shapes.add_textbox(
            Inches(7), Inches(6), Inches(2.5), Inches(0.5)
        )
        date_frame = date_shape.text_frame
        date_frame.text = f"Généré le: {pd.Timestamp.now().strftime('%d/%m/%Y')}"
        date_frame.paragraphs[0].font.size = Pt(10)
        date_frame.paragraphs[0].font.color.rgb = self.colors['dark']
    
    def _create_executive_summary_slide(self, prs: Presentation, data: Dict):
        """Crée le résumé exécutif"""
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        
        title = slide.shapes.title
        title.text = "Résumé Exécutif"
        title.text_frame.paragraphs[0].font.color.rgb = self.colors['primary']
        
        # Calcul des KPIs principaux
        monthly_data = data.get('monthly_data', pd.DataFrame())
        agents_data = data.get('agents_data', pd.DataFrame())
        
        kpis = self._calculate_executive_kpis(monthly_data, agents_data)
        
        # Créer les zones de texte pour les KPIs
        self._add_kpi_boxes(slide, kpis)
        
        # Points clés
        content_box = slide.shapes.add_textbox(
            Inches(1), Inches(3.5), Inches(8), Inches(3)
        )
        
        tf = content_box.text_frame
        tf.text = "Points Clés:"
        
        key_points = self._generate_key_points(kpis, monthly_data)
        for point in key_points:
            p = tf.add_paragraph()
            p.text = f"• {point}"
            p.level = 1
            p.font.size = Pt(12)
    
    def _create_monthly_analysis_slide(self, prs: Presentation, data: Dict, figures: Dict):
        """Crée l'analyse mensuelle"""
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Titre
        title_shape = slide.shapes.add_textbox(
            Inches(1), Inches(0.5), Inches(8), Inches(0.8)
        )
        title_frame = title_shape.text_frame
        title_frame.text = "Analyse Mensuelle des Performances"
        title_frame.paragraphs[0].font.size = Pt(24)
        title_frame.paragraphs[0].font.color.rgb = self.colors['primary']
        title_frame.paragraphs[0].font.bold = True
        
        # Graphique principal (si disponible)
        if 'volume_calls' in figures:
            self._add_chart_to_slide(slide, figures['volume_calls'], 
                                   Inches(1), Inches(1.5), Inches(8), Inches(4))
        
        # Tableau de données
        monthly_data = data.get('monthly_data', pd.DataFrame())
        if not monthly_data.empty:
            self._add_data_table(slide, monthly_data, 
                                Inches(1), Inches(6), Inches(8), Inches(1.5))
    
    def _create_agents_performance_slide(self, prs: Presentation, data: Dict, figures: Dict):
        """Crée l'analyse des agents"""
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Titre
        title_shape = slide.shapes.add_textbox(
            Inches(1), Inches(0.5), Inches(8), Inches(0.8)
        )
        title_frame = title_shape.text_frame
        title_frame.text = "Performance Individuelle des Agents"
        title_frame.paragraphs[0].font.size = Pt(24)
        title_frame.paragraphs[0].font.color.rgb = self.colors['primary']
        title_frame.paragraphs[0].font.bold = True
        
        # Données des agents
        agents_data = data.get('agents_data', pd.DataFrame())
        if not agents_data.empty:
            # Top performer
            if 'appels_traites' in agents_data and 'agent' in agents_data:
                top_agent = agents_data.loc[agents_data['appels_traites'].idxmax(), 'agent']
                
                highlight_box = slide.shapes.add_textbox(
                    Inches(6), Inches(1.5), Inches(3), Inches(1.5)
                )
                hf = highlight_box.text_frame
                hf.text = f"Top Performer\n{top_agent}"
                hf.paragraphs[0].font.size = Pt(16)
                hf.paragraphs[0].font.bold = True
                hf.paragraphs[0].font.color.rgb = self.colors['secondary']
            
            # Tableau des agents
            self._add_agents_table(slide, agents_data, 
                                 Inches(1), Inches(3), Inches(8), Inches(3))
    
    def _create_kpi_trends_slide(self, prs: Presentation, data: Dict, figures: Dict):
        """Crée l'analyse des tendances et KPIs"""
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Titre
        title_shape = slide.shapes.add_textbox(
            Inches(1), Inches(0.5), Inches(8), Inches(0.8)
        )
        title_frame = title_shape.text_frame
        title_frame.text = "Indicateurs Clés et Tendances"
        title_frame.paragraphs[0].font.size = Pt(24)
        title_frame.paragraphs[0].font.color.rgb = self.colors['primary']
        title_frame.paragraphs[0].font.bold = True
        
        # KPIs en colonnes
        monthly_data = data.get('monthly_data', pd.DataFrame())
        if not monthly_data.empty:
            kpis = self._calculate_executive_kpis(monthly_data, data.get('agents_data', pd.DataFrame()))
            self._create_kpi_grid(slide, kpis)
    
    def _create_resolution_analysis_slide(self, prs: Presentation, data: Dict):
        """Analyse de la résolution des appels"""
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        
        title = slide.shapes.title
        title.text = "Analyse de la Résolution"
        title.text_frame.paragraphs[0].font.color.rgb = self.colors['primary']
        
        content = slide.placeholders[1]
        tf = content.text_frame
        
        # Calculer les métriques de résolution
        resolution_data = data.get('resolution_data', pd.DataFrame())
        monthly_data = data.get('monthly_data', pd.DataFrame())
        
        if not monthly_data.empty and 'appels_presentes' in monthly_data and 'appels_traites' in monthly_data:
            taux_global = (monthly_data['appels_traites'].sum() / 
                          monthly_data['appels_presentes'].sum() * 100)
            
            tf.text = f"Taux de Résolution Global: {taux_global:.1f}%"
            
            # Analyse par période
            if len(monthly_data) > 1:
                evolution = self._calculate_resolution_evolution(monthly_data)
                for point in evolution:
                    p = tf.add_paragraph()
                    p.text = f"• {point}"
                    p.level = 1
    
    def _create_recommendations_slide(self, prs: Presentation, data: Dict):
        """Crée les recommandations"""
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        
        title = slide.shapes.title
        title.text = "Recommandations"
        title.text_frame.paragraphs[0].font.color.rgb = self.colors['primary']
        
        content = slide.placeholders[1]
        tf = content.text_frame
        
        recommendations = self._generate_recommendations(data)
        tf.text = "Actions Recommandées:"
        
        for rec in recommendations:
            p = tf.add_paragraph()
            p.text = f"• {rec}"
            p.level = 1
            p.font.size = Pt(12)
    
    def _calculate_executive_kpis(self, monthly_data: pd.DataFrame, agents_data: pd.DataFrame) -> Dict:
        """Calcule les KPIs pour le résumé exécutif"""
        kpis = {}
        
        if not monthly_data.empty:
            kpis['total_volume'] = monthly_data.get('appels_presentes', pd.Series([0])).sum()
            kpis['total_traites'] = monthly_data.get('appels_traites', pd.Series([0])).sum()
            kpis['taux_resolution'] = (kpis['total_traites'] / kpis['total_volume'] * 100) if kpis['total_volume'] > 0 else 0
            kpis['duree_moyenne'] = monthly_data.get('duree_moyenne_conv', pd.Series([0])).mean()
            kpis['periode_couverte'] = len(monthly_data)
        
        if not agents_data.empty:
            kpis['nb_agents'] = len(agents_data)
            kpis['agent_top'] = agents_data.loc[agents_data.get('appels_traites', pd.Series([0])).idxmax(), 'agent'] if 'agent' in agents_data and 'appels_traites' in agents_data else 'N/A'
        
        return kpis
    
    def _add_kpi_boxes(self, slide, kpis: Dict):
        """Ajoute les boîtes KPI à la slide"""
        box_width = Inches(1.8)
        box_height = Inches(1.2)
        start_x = Inches(1)
        start_y = Inches(1.5)
        spacing = Inches(0.2)
        
        kpi_items = [
            ("Volume Total", f"{kpis.get('total_volume', 0):,.0f}", self.colors['primary']),
            ("Taux Résolution", f"{kpis.get('taux_resolution', 0):.1f}%", self.colors['secondary']),
            ("Durée Moyenne", f"{kpis.get('duree_moyenne', 0):.1f} min", self.colors['warning']),
            ("Agents Actifs", f"{kpis.get('nb_agents', 0)}", self.colors['accent'])
        ]
        
        for i, (label, value, color) in enumerate(kpi_items):
            x_pos = start_x + i * (box_width + spacing)
            
            # Boîte de fond
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, x_pos, start_y, box_width, box_height
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = self.colors['light']
            shape.line.color.rgb = color
            shape.line.width = Pt(2)
            
            # Texte de la valeur
            value_box = slide.shapes.add_textbox(
                x_pos, start_y + Inches(0.1), box_width, Inches(0.6)
            )
            vf = value_box.text_frame
            vf.text = value
            vf.paragraphs[0].alignment = PP_ALIGN.CENTER
            vf.paragraphs[0].font.size = Pt(18)
            vf.paragraphs[0].font.bold = True
            vf.paragraphs[0].font.color.rgb = color
            
            # Texte du label
            label_box = slide.shapes.add_textbox(
                x_pos, start_y + Inches(0.7), box_width, Inches(0.4)
            )
            lf = label_box.text_frame
            lf.text = label
            lf.paragraphs[0].alignment = PP_ALIGN.CENTER
            lf.paragraphs[0].font.size = Pt(10)
            lf.paragraphs[0].font.color.rgb = self.colors['dark']
    
    def _add_chart_to_slide(self, slide, figure: go.Figure, left, top, width, height):
        """Ajoute un graphique Plotly à la slide (version simplifiée)"""
        # Note: Pour une implémentation complète, il faudrait convertir le graphique en image
        # et l'insérer dans la slide. Ici, on crée un placeholder.
        
        chart_box = slide.shapes.add_textbox(left, top, width, height)
        cf = chart_box.text_frame
        cf.text = "Graphique des performances mensuelles\n(Visualisation interactive disponible dans l'application)"
        cf.paragraphs[0].alignment = PP_ALIGN.CENTER
        cf.paragraphs[0].font.size = Pt(14)
        cf.paragraphs[0].font.color.rgb = self.colors['dark']
    
    def _add_data_table(self, slide, data: pd.DataFrame, left, top, width, height):
        """Ajoute un tableau de données"""
        # Créer un tableau simple avec les données principales
        if data.empty:
            return
        
        table_data = []
        if 'mois' in data:
            table_data.append(['Mois'] + data['mois'].tolist())
        if 'appels_traites' in data:
            table_data.append(['Appels Traités'] + [str(x) for x in data['appels_traites'].tolist()])
        if 'appels_presentes' in data:
            table_data.append(['Appels Présentés'] + [str(x) for x in data['appels_presentes'].tolist()])
        
        # Créer une zone de texte formatée comme tableau
        if table_data:
            table_text = '\n'.join(['\t'.join(row) for row in table_data])
            table_box = slide.shapes.add_textbox(left, top, width, height)
            tf = table_box.text_frame
            tf.text = table_text
            tf.paragraphs[0].font.size = Pt(10)
            tf.paragraphs[0].font.name = 'Courier New'  # Police monospace pour alignement
    
    def _add_agents_table(self, slide, agents_data: pd.DataFrame, left, top, width, height):
        """Ajoute le tableau des agents"""
        if agents_data.empty or 'agent' not in agents_data:
            return
        
        agents_text = "Performance des Agents:\n\n"
        for _, row in agents_data.iterrows():
            agent_name = row['agent']
            appels = row.get('appels_traites', 'N/A')
            agents_text += f"• {agent_name}: {appels} appels traités\n"
        
        table_box = slide.shapes.add_textbox(left, top, width, height)
        tf = table_box.text_frame
        tf.text = agents_text
        tf.paragraphs[0].font.size = Pt(12)
    
    def _create_kpi_grid(self, slide, kpis: Dict):
        """Crée une grille de KPIs"""
        # Version simplifiée - utilise des zones de texte
        y_pos = Inches(2)
        
        kpi_text = f"""INDICATEURS CLÉS DE PERFORMANCE
        
• Volume Total Traité: {kpis.get('total_volume', 0):,} appels
• Taux de Résolution: {kpis.get('taux_resolution', 0):.1f}%
• Durée Moyenne: {kpis.get('duree_moyenne', 0):.1f} minutes
• Période Analysée: {kpis.get('periode_couverte', 0)} mois
• Agents Actifs: {kpis.get('nb_agents', 0)} agents
        """
        
        grid_box = slide.shapes.add_textbox(Inches(1), y_pos, Inches(8), Inches(3))
        gf = grid_box.text_frame
        gf.text = kpi_text
        gf.paragraphs[0].font.size = Pt(14)
        gf.paragraphs[0].font.bold = True
        gf.paragraphs[0].font.color.rgb = self.colors['primary']
        
        # Ajouter des puces pour les autres paragraphes
        for i in range(1, len(gf.paragraphs)):
            gf.paragraphs[i].font.size = Pt(12)
    
    def _generate_key_points(self, kpis: Dict, monthly_data: pd.DataFrame) -> List[str]:
        """Génère les points clés pour le résumé exécutif"""
        points = []
        
        # Analyse du volume
        if kpis.get('total_volume', 0) > 0:
            points.append(f"Volume total de {kpis['total_volume']:,} appels traités sur la période")
        
        # Analyse du taux de résolution
        taux = kpis.get('taux_resolution', 0)
        if taux >= 95:
            points.append(f"Excellent taux de résolution de {taux:.1f}% (objectif dépassé)")
        elif taux >= 90:
            points.append(f"Bon taux de résolution de {taux:.1f}% (proche de l'objectif)")
        else:
            points.append(f"Taux de résolution de {taux:.1f}% nécessite une attention")
        
        # Analyse de l'efficacité
        duree = kpis.get('duree_moyenne', 0)
        if duree <= 5:
            points.append("Durée moyenne de conversation optimale (≤5 min)")
        else:
            points.append(f"Durée moyenne de {duree:.1f} min peut être optimisée")
        
        # Analyse des tendances
        if not monthly_data.empty and len(monthly_data) > 1:
            if 'appels_traites' in monthly_data:
                evolution = self._calculate_volume_trend(monthly_data['appels_traites'])
                if evolution > 5:
                    points.append("Tendance positive du volume d'appels (+{:.1f}%)".format(evolution))
                elif evolution < -5:
                    points.append("Tendance négative du volume d'appels ({:.1f}%)".format(evolution))
                else:
                    points.append("Volume d'appels stable sur la période")
        
        return points
    
    def _calculate_volume_trend(self, volume_series: pd.Series) -> float:
        """Calcule la tendance d'évolution du volume"""
        if len(volume_series) < 2:
            return 0
        
        first_half = volume_series[:len(volume_series)//2].mean()
        second_half = volume_series[len(volume_series)//2:].mean()
        
        if first_half == 0:
            return 0
        
        return ((second_half - first_half) / first_half) * 100
    
    def _calculate_resolution_evolution(self, monthly_data: pd.DataFrame) -> List[str]:
        """Calcule l'évolution de la résolution"""
        evolution_points = []
        
        if 'appels_presentes' in monthly_data and 'appels_traites' in monthly_data:
            # Taux de résolution par mois
            taux_mensuel = (monthly_data['appels_traites'] / monthly_data['appels_presentes'] * 100)
            
            # Meilleur et pire mois
            best_month_idx = taux_mensuel.idxmax()
            worst_month_idx = taux_mensuel.idxmin()
            
            if 'mois' in monthly_data:
                best_month = monthly_data.loc[best_month_idx, 'mois']
                worst_month = monthly_data.loc[worst_month_idx, 'mois']
                
                evolution_points.append(
                    f"Meilleur mois: {best_month} ({taux_mensuel[best_month_idx]:.1f}%)"
                )
                evolution_points.append(
                    f"Mois le plus difficile: {worst_month} ({taux_mensuel[worst_month_idx]:.1f}%)"
                )
            
            # Tendance générale
            if len(taux_mensuel) > 2:
                trend = self._calculate_volume_trend(taux_mensuel)
                if trend > 1:
                    evolution_points.append("Tendance d'amélioration continue")
                elif trend < -1:
                    evolution_points.append("Attention: tendance de dégradation")
                else:
                    evolution_points.append("Performance stable")
        
        return evolution_points
    
    def _generate_recommendations(self, data: Dict) -> List[str]:
        """Génère des recommandations basées sur l'analyse"""
        recommendations = []
        
        monthly_data = data.get('monthly_data', pd.DataFrame())
        agents_data = data.get('agents_data', pd.DataFrame())
        
        # Recommandations basées sur le taux de résolution
        if not monthly_data.empty and 'appels_presentes' in monthly_data and 'appels_traites' in monthly_data:
            taux_global = (monthly_data['appels_traites'].sum() / 
                          monthly_data['appels_presentes'].sum() * 100)
            
            if taux_global < 90:
                recommendations.append("Améliorer le taux de résolution par formation complémentaire")
                recommendations.append("Analyser les causes de non-résolution des appels")
            
            if taux_global > 98:
                recommendations.append("Optimiser les processus pour maintenir l'excellence")
        
        # Recommandations basées sur la durée
        if not monthly_data.empty and 'duree_moyenne_conv' in monthly_data:
            duree_moy = monthly_data['duree_moyenne_conv'].mean()
            
            if duree_moy > 6:
                recommendations.append("Réduire la durée moyenne par optimisation des scripts")
                recommendations.append("Former les agents aux techniques de communication efficace")
            elif duree_moy < 3:
                recommendations.append("Vérifier la qualité du service malgré la rapidité")
        
        # Recommandations sur les ressources
        if not monthly_data.empty and 'nb_agents_max' in monthly_data:
            variation_agents = monthly_data['nb_agents_max'].std()
            
            if variation_agents > 1:
                recommendations.append("Stabiliser l'effectif pour une meilleure prévisibilité")
        
        # Recommandations basées sur les agents
        if not agents_data.empty and 'appels_traites' in agents_data:
            # Disparité entre agents
            if len(agents_data) > 1:
                cv = agents_data['appels_traites'].std() / agents_data['appels_traites'].mean()
                if cv > 0.5:
                    recommendations.append("Équilibrer la charge de travail entre agents")
                    recommendations.append("Identifier et partager les bonnes pratiques")
        
        # Recommandations générales si aucune spécifique
        if not recommendations:
            recommendations.extend([
                "Maintenir le niveau de performance actuel",
                "Continuer le monitoring régulier des indicateurs",
                "Planifier des sessions de formation continue"
            ])
        
        return recommendations
    
    def save_presentation(self, prs: Presentation, filename: str = None) -> str:
        """Sauvegarde la présentation et retourne le chemin"""
        if not filename:
            filename = f"rapport_telephonie_{pd.Timestamp.now().strftime('%Y%m%d_%H%M%S')}.pptx"
        
        # Utiliser un répertoire temporaire pour le déploiement web
        temp_dir = tempfile.gettempdir()
        filepath = os.path.join(temp_dir, filename)
        
        prs.save(filepath)
        return filepath
    
    def get_presentation_bytes(self, prs: Presentation) -> bytes:
        """Retourne la présentation sous forme de bytes pour téléchargement"""
        temp_buffer = io.BytesIO()
        prs.save(temp_buffer)
        temp_buffer.seek(0)
        return temp_buffer.getvalue()

# Fonctions utilitaires
def create_sample_presentation() -> Presentation:
    """Crée une présentation d'exemple"""
    generator = PowerPointReportGenerator()
    
    # Données d'exemple
    sample_data = {
        'monthly_data': pd.DataFrame({
            'mois': ['Janvier', 'Février', 'Mars', 'Avril'],
            'appels_presentes': [594, 554, 584, 641],
            'appels_traites': [570, 543, 550, 626],
            'duree_moyenne_conv': [5.51, 5.06, 5.14, 5.26],
            'nb_agents_max': [3, 3, 2, 2]
        }),
        'agents_data': pd.DataFrame({
            'agent': ['Fabienne Cocquart', 'Philippe Kubler'],
            'appels_presentes': [594, 570],
            'appels_traites': [570, 543]
        })
    }
    
    return generator.create_presentation(sample_data, {})

def format_number_french(number: float) -> str:
    """Formate un nombre selon les conventions françaises"""
    return f"{number:,.0f}".replace(',', ' ')

def calculate_performance_score(taux_resolution: float, duree_moyenne: float, volume: int) -> float:
    """Calcule un score de performance global"""
    # Score basé sur différents critères (0-100)
    score_resolution = min(taux_resolution, 100)
    score_efficacite = max(0, 100 - (duree_moyenne - 3) * 10)  # Optimal à 3min
    score_volume = min(volume / 10, 100)  # Score relatif au volume
    
    # Moyenne pondérée
    score_global = (score_resolution * 0.5 + score_efficacite * 0.3 + score_volume * 0.2)
    return max(0, min(100, score_global))

if __name__ == "__main__":
    # Test de création d'une présentation
    sample_prs = create_sample_presentation()
    print("Présentation d'exemple créée avec succès")
